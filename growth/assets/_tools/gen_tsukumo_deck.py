#!/usr/bin/env python3
"""Generate the tsukumo PPTX master template (acid/brutalist) via python-pptx.
5 master/example slides: title / section-divider / content / case-study / contact.
Copy is PLACEHOLDER/qualitative — content-lead fills real words; NO fabricated proof.
Run with a venv that has python-pptx. Writes growth/agency/brand/tsukumo-master.pptx.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK   = RGBColor(0x12, 0x12, 0x12)
INK2  = RGBColor(0x1b, 0x1a, 0x18)
BONE  = RGBColor(0xf3, 0xf1, 0xea)
DIM   = RGBColor(0xb8, 0xb3, 0xa8)
ACID  = RGBColor(0xc8, 0xff, 0x00)
ACIDK = RGBColor(0x0e, 0x11, 0x00)
DISP  = "Archivo"
MONO  = "Space Mono"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
OUT = os.path.join(os.path.dirname(__file__), "..", "..", "agency", "brand", "tsukumo-master.pptx")

def _bg(slide, color=INK):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    s.shadow.inherit = False
    slide.shapes._spTree.remove(s._element); slide.shapes._spTree.insert(2, s._element)
    return s

def _rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    s.shadow.inherit = False
    return s

def _text(slide, x, y, w, h, runs, size, font=DISP, color=BONE, bold=True,
          align=PP_ALIGN.LEFT, spacing=-0.03, line=0.95):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line
        r = p.add_run(); r.text = txt
        f = r.font; f.size = Pt(size); f.name = font; f.bold = bold
        f.color.rgb = col
    return tb

def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    _rect(s, 0, 0, EMU_W, Inches(0.09), ACID)
    _rect(s, Inches(0.8), Inches(0.78), Inches(0.3), Inches(0.3), ACID)
    _text(s, Inches(1.2), Inches(0.66), Inches(6), Inches(0.6),
          [("tsukumo", BONE)], 24, DISP, BONE, True, spacing=-0.02)
    _text(s, Inches(8.6), Inches(0.74), Inches(3.9), Inches(0.5),
          [("a dev studio + AI consulting", DIM)], 13, MONO, DIM, False, PP_ALIGN.RIGHT)
    _text(s, Inches(0.72), Inches(2.5), Inches(12), Inches(2.6),
          [("run AI in", BONE), ("production.", ACID)], 88, DISP, BONE, True, line=0.9)
    _rect(s, Inches(0.8), Inches(5.75), Inches(1.1), Inches(0.10), ACID)
    _text(s, Inches(0.8), Inches(6.0), Inches(11), Inches(0.8),
          [("agentic operators for your dev team — augment, never replace.", BONE)],
          16, MONO, BONE, False, line=1.2)
    _text(s, Inches(0.8), Inches(6.95), Inches(6), Inches(0.4),
          [("tsukumo.ch", ACID)], 13, MONO, ACID, False)
    return s

def divider_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    _rect(s, 0, 0, EMU_W, Inches(0.09), ACID)
    _text(s, Inches(0.8), Inches(0.7), Inches(4), Inches(0.5),
          [("section · 0X", ACID)], 14, MONO, ACID, False)
    _text(s, Inches(0.72), Inches(2.8), Inches(12), Inches(2),
          [("approach", BONE)], 110, DISP, BONE, True, line=0.9)
    _rect(s, Inches(0.8), Inches(5.2), Inches(1.4), Inches(0.12), ACID)
    return s

def content_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    _rect(s, 0, 0, EMU_W, Inches(0.09), ACID)
    _text(s, Inches(0.8), Inches(0.7), Inches(6), Inches(0.5),
          [("01 — content", ACID)], 14, MONO, ACID, False)
    _text(s, Inches(0.72), Inches(1.3), Inches(11.8), Inches(1.4),
          [("a section title goes here", BONE)], 44, DISP, BONE, True, line=0.95)
    _rect(s, Inches(0.8), Inches(2.95), Inches(11.7), Pt(1.4), RGBColor(0x39,0x37,0x33))
    _text(s, Inches(0.8), Inches(3.2), Inches(11.5), Inches(3.4),
          [("— point one, kept short and declarative", BONE),
           ("— point two, qualitative (no fabricated numbers)", BONE),
           ("— point three, one idea per line", BONE)],
          18, MONO, BONE, False, line=1.5)
    _text(s, Inches(0.8), Inches(6.95), Inches(6), Inches(0.4),
          [("tsukumo", ACID)], 12, MONO, ACID, False)
    return s

def case_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    _rect(s, 0, 0, EMU_W, Inches(0.09), ACID)
    _text(s, Inches(0.8), Inches(0.7), Inches(6), Inches(0.5),
          [("case · studio", ACID)], 14, MONO, ACID, False)
    _text(s, Inches(11.0), Inches(0.5), Inches(1.6), Inches(1.2),
          [("01", INK2)], 64, DISP, INK2, True, PP_ALIGN.RIGHT)
    _text(s, Inches(0.72), Inches(2.3), Inches(12), Inches(1.6),
          [("our own stack", BONE)], 96, DISP, BONE, True, line=0.9)
    _rect(s, Inches(0.8), Inches(4.2), Inches(1.1), Inches(0.10), ACID)
    _text(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.6),
          [("we run agents in production: WRAI.TH (orchestration),", BONE),
           ("trovex (context), yoru (observability) — open source.", BONE)],
          20, MONO, BONE, False, line=1.4)
    _text(s, Inches(0.8), Inches(6.95), Inches(6), Inches(0.4),
          [("the studio is the proof", DIM)], 13, MONO, DIM, False)
    return s

def contact_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    _rect(s, 0, 0, EMU_W, Inches(0.09), ACID)
    _text(s, Inches(0.8), Inches(0.7), Inches(4), Inches(0.5),
          [("contact", ACID)], 14, MONO, ACID, False)
    _text(s, Inches(0.72), Inches(2.7), Inches(12), Inches(2),
          [("let's talk.", ACID)], 120, DISP, ACID, True, line=0.9)
    _text(s, Inches(0.8), Inches(5.3), Inches(11), Inches(1),
          [("tsukumo.ch  ·  [ contact email ]", BONE)], 22, MONO, BONE, False, line=1.3)
    _text(s, Inches(0.8), Inches(6.95), Inches(8), Inches(0.4),
          [("a dev studio + AI consulting · augment, never replace", DIM)], 12, MONO, DIM, False)
    return s

def main():
    prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
    title_slide(prs); divider_slide(prs); content_slide(prs); case_slide(prs); contact_slide(prs)
    prs.save(os.path.abspath(OUT))
    print("wrote", os.path.abspath(OUT), "—", len(prs.slides._sldIdLst), "slides")

if __name__ == "__main__":
    main()
